import base64
import datetime
import tempfile
from PIL import Image
from PyPDF2 import PdfReader
import concurrent
from concurrent.futures import ThreadPoolExecutor
from pdf2image import convert_from_path
from rest_framework.views import APIView
from rest_framework.response import Response
from rest_framework import status
from pptx import Presentation
from fitz import open as open_pdf
import os
import io
import pytesseract
import requests
from rest_framework.decorators import api_view
from azure.cognitiveservices.vision.computervision import ComputerVisionClient
from azure.cognitiveservices.vision.computervision.models import OperationStatusCodes
from msrest.authentication import CognitiveServicesCredentials
import subprocess
from rest_framework.parsers import MultiPartParser
from api import models, serializers
from decouple import config
import threading

# Global image description count and last recharge time
image_description_count = 20
last_recharge_time = datetime.datetime.now()

# Tesseract OCR configuration
pytesseract.pytesseract.tesseract_cmd = 'C:\\Program Files\\Tesseract-OCR\\tesseract.exe' if os.name == 'nt' else '/usr/bin/tesseract'


# Azure Cognitive Services configuration
# subscription_key = os.getenv("AZURE_SUBSCRIPTION_KEY", "default")
subscription_key = config('AZURE_SUBSCRIPTION_KEY', default='default')
endpoint = os.getenv("AZURE_ENDPOINT", "https://scribemeocr.cognitiveservices.azure.com/")
computervision_client = ComputerVisionClient(endpoint, CognitiveServicesCredentials(subscription_key))


def describe_image_with_gpt(base64_image, prompt_text="Describe this image"):
    # api_key = os.getenv("OPENAI_API_KEY")  # Ensure this is set in your environment variables
    api_key = config('OPENAI_API_KEY', default='default')  # Ensure this is set in your environment variables
    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Bearer {api_key}",
    }

    payload = {
        "model": "gpt-4o-2024-08-06",
        "messages": [
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": prompt_text},
                    {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{base64_image}"}}
                ]
            }
        ],
        "max_tokens": 325
    }

    response = requests.post("https://api.openai.com/v1/chat/completions", headers=headers, json=payload)
    print(response)
    response.raise_for_status()
    response_json = response.json()
    return response_json["choices"][0]["message"]["content"]


# class DescribeImageView(APIView):
#     def post(self, request):
#         image_file = request.FILES.get("image")
#         language = request.data.get("language", "English")

#         if not image_file:
#             return Response({"error": "Image file is required."}, status=status.HTTP_400_BAD_REQUEST)

#         try:
#             image = Image.open(image_file)
#             buffered = io.BytesIO()
#             image.save(buffered, format="JPEG")
#             base64_image = base64.b64encode(buffered.getvalue()).decode("utf-8")

#             prompt_texts = {
#                 "English": "Describe this image in detail.",
#                 "Arabic": "صف هذه الصورة بالتفصيل.",
#                 "Spanish": "Describe esta imagen en detalle."
#             }
#             prompt_text = prompt_texts.get(language, "Describe this image in detail.")

#             description = describe_image_with_gpt(base64_image, prompt_text)
#             return Response({"description": description}, status=status.HTTP_200_OK)

#         except Exception as e:
#             return Response({"error": str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

class DescribeImageView(APIView):
    def post(self, request):
        image_files = request.FILES.getlist("images")  # Get multiple images
        language = request.data.get("language", "English")
        print(image_files)
        if not image_files:
            return Response({"error": "At least one image file is required."}, status=status.HTTP_400_BAD_REQUEST)

        prompt_texts = {
            "English": "Describe this image in detail.",
            "Arabic": "صف هذه الصورة بالتفصيل.",
            "Spanish": "Describe esta imagen en detalle."
        }
        prompt_text = prompt_texts.get(language, "Describe this image in detail.")

        descriptions = []

        try:
            for image_file in image_files:
                image = Image.open(image_file)
                buffered = io.BytesIO()
                image.save(buffered, format="JPEG")
                base64_image = base64.b64encode(buffered.getvalue()).decode("utf-8")

                description = describe_image_with_gpt(base64_image, prompt_text)
                descriptions.append({
                    "filename": image_file.name,
                    "description": description
                })

            return Response({"descriptions": descriptions}, status=status.HTTP_200_OK)

        except Exception as e:
            return Response({"error": str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

def perform_ocr(image_path, lang="eng"):
    image = Image.open(image_path)
    return pytesseract.image_to_string(image, lang=lang)


def analyze_image_with_ocr_with_arabic(image_path):
    try:
        with open(image_path, "rb") as image_stream:
            ocr_result = computervision_client.recognize_printed_text_in_stream(image=image_stream, language="ar")
        return "\n".join(" ".join(word.text for word in line.words) for region in ocr_result.regions for line in region.lines)
    except Exception as e:
        return f"Error: {e}"


# class ExtractTextFromPDFView(APIView):
#     def post(self, request):
#         pdf_file = request.FILES.get("pdf_file")
#         ocr_option = request.data.get("ocr", False)
#         remaining_images = int(request.data.get("rImages")) if request.data.get("rImages") else 25
        
#         image_description_option = request.data.get("image_description", False)
#         language = request.data.get("language", "English")

#         if not pdf_file:
#             return Response({"error": "PDF file is required."}, status=status.HTTP_400_BAD_REQUEST)

#         with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as temp_pdf:
#             temp_pdf.write(pdf_file.read())
#             temp_pdf_path = temp_pdf.name

#         try:
#             text_content = ""
#             image_description_count = 0

#             with open_pdf(temp_pdf_path) as pdf_document:
#                 for page_number, page in enumerate(pdf_document, start=1):
#                     text_content += f"Page {page_number}:\n{page.get_text('text')}\n"

#                     if str(ocr_option) == 'true' or str(image_description_option) == 'true':
#                         for img in page.get_images(full=True):
#                             xref = img[0]
#                             image_bytes = pdf_document.extract_image(xref)["image"]

#                             with tempfile.NamedTemporaryFile(delete=False, suffix=".jpg") as temp_image:
#                                 temp_image.write(image_bytes)
#                                 temp_image_path = temp_image.name

#                             if str(ocr_option) == 'true':
#                                 lang_map = {"English": "eng", "Spanish": "spa", "Arabic": "ara"}
#                                 lang_code = lang_map.get(language, "eng")
#                                 text_content += f"\n OCR Text from image on page {page_number}: {perform_ocr(temp_image_path, lang_code)}\n"

#                             if str(image_description_option) == 'true' and remaining_images > 0:
#                                 prompt_texts = {
#                                     "English": "Describe this image in detail.",
#                                     "Arabic": "صف هذه الصورة بالتفصيل.",
#                                     "Spanish": "Describe esta imagen en detalle."
#                                 }
#                                 prompt_text = prompt_texts.get(language, "Describe this image in detail.")
#                                 gpt_description = describe_image_with_gpt(
#                                     base64.b64encode(image_bytes).decode("utf-8"),
#                                     prompt_text
#                                 )
#                                 text_content += f"\n Image description on page {page_number}: {gpt_description}\n"
#                                 image_description_count += 1
#                                 remaining_images -= 1

#                             os.remove(temp_image_path)

#             os.remove(temp_pdf_path)
#             return Response({"text_content": text_content, "count": image_description_count}, status=status.HTTP_200_OK)

#         except Exception as e:
#             if os.path.exists(temp_pdf_path):
#                 os.remove(temp_pdf_path)
#             return Response({"error": str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

# def process_page(page, page_number, ocr_option, image_description_option, language, remaining_images):
#     """Processes a single page: extracts text, performs OCR, and describes images if required."""
#     text_content = f"Page {page_number}:\n{page.get_text('text')}\n"
#     image_description_count = 0

#     if ocr_option or image_description_option:
#         for img in page.get_images(full=True):
#             xref = img[0]
#             image_bytes = page.parent.extract_image(xref)["image"]

#             with tempfile.NamedTemporaryFile(delete=False, suffix=".jpg") as temp_image:
#                 temp_image.write(image_bytes)
#                 temp_image_path = temp_image.name

#             if ocr_option:
#                 lang_map = {"English": "eng", "Spanish": "spa", "Arabic": "ara"}
#                 lang_code = lang_map.get(language, "eng")
#                 text_content += f"\n OCR Text from image on page {page_number}: {perform_ocr(temp_image_path, lang_code)}\n"

#             if image_description_option and remaining_images > 0:
#                 prompt_texts = {
#                     "English": "Describe this image in detail.",
#                     "Arabic": "صف هذه الصورة بالتفصيل.",
#                     "Spanish": "Describe esta imagen en detalle."
#                 }
#                 prompt_text = prompt_texts.get(language, "Describe this image in detail.")
#                 gpt_description = describe_image_with_gpt(
#                     base64.b64encode(image_bytes).decode("utf-8"),
#                     prompt_text
#                 )
#                 text_content += f"\n Image description on page {page_number}: {gpt_description}\n"
#                 image_description_count += 1
#                 remaining_images -= 1

#             os.remove(temp_image_path)

#     return text_content, image_description_count

# class ExtractTextFromPDFView(APIView):
#     def post(self, request):
#         pdf_file = request.FILES.get("pdf_file")
#         ocr_option = str(request.data.get("ocr", "false")).lower() == "true"
#         image_description_option = str(request.data.get("image_description", "false")).lower() == "true"
#         remaining_images = int(request.data.get("rImages", 25))
#         language = request.data.get("language", "English")
#         batch_size = 5  

#         if not pdf_file:
#             return Response({"error": "PDF file is required."}, status=status.HTTP_400_BAD_REQUEST)

#         with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as temp_pdf:
#             temp_pdf.write(pdf_file.read())
#             temp_pdf_path = temp_pdf.name

#         try:
#             text_content = ""
#             image_description_count = 0

#             with open_pdf(temp_pdf_path) as pdf_document:
#                 num_pages = len(pdf_document)
#                 with ThreadPoolExecutor() as executor:
#                     futures = []
#                     for i in range(0, num_pages, batch_size):
#                         batch_pages = [pdf_document[j] for j in range(i, min(i + batch_size, num_pages))]
#                         for page_number, page in enumerate(batch_pages, start=i + 1):
#                             futures.append(executor.submit(
#                                 process_page, page, page_number, ocr_option, image_description_option, language, remaining_images
#                             ))

#                     for future in futures:
#                         page_text, img_count = future.result()
#                         text_content += page_text
#                         image_description_count += img_count
#                         remaining_images -= img_count  # Update remaining images count

#             os.remove(temp_pdf_path)
#             return Response({"text_content": text_content, "count": image_description_count}, status=status.HTTP_200_OK)

#         except Exception as e:
#             if os.path.exists(temp_pdf_path):
#                 os.remove(temp_pdf_path)
#             return Response({"error": str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

class ThreadSafeCounter:
    def __init__(self, initial_value):
        self.value = initial_value
        self.lock = threading.Lock()

    def decrement(self):
        with self.lock:
            if self.value > 0:
                self.value -= 1
                return True
            return False
        
    def get_value(self):
        with self.lock:
            return self.value

def process_page(page, page_number, ocr_option, image_description_option, language, remaining_images_counter):
    text_content = f"Page {page_number}:\n{page.get_text('text')}\n"
    image_description_count = 0

    if ocr_option or image_description_option:
        for img in page.get_images(full=True):
            xref = img[0]
            image_bytes = page.parent.extract_image(xref)["image"]

            with tempfile.NamedTemporaryFile(delete=False, suffix=".jpg") as temp_image:
                temp_image.write(image_bytes)
                temp_image_path = temp_image.name

            if ocr_option:
                lang_map = {"English": "eng", "Spanish": "spa", "Arabic": "ara"}
                lang_code = lang_map.get(language, "eng")
                text_content += f"\n OCR Text from image on page {page_number}: {perform_ocr(temp_image_path, lang_code)}\n"

            if image_description_option and remaining_images_counter.decrement():
                prompt_texts = {
                    "English": "Describe this image in detail.",
                    "Arabic": "صف هذه الصورة بالتفصيل.",
                    "Spanish": "Describe esta imagen en detalle."
                }
                prompt_text = prompt_texts.get(language, "Describe this image in detail.")
                gpt_description = describe_image_with_gpt(
                    base64.b64encode(image_bytes).decode("utf-8"),
                    prompt_text
                )
                text_content += f"\n Image description on page {page_number}: {gpt_description}\n"
                image_description_count += 1

            os.remove(temp_image_path)

    return text_content, image_description_count

class ExtractTextFromPDFView(APIView):
    def post(self, request):
        pdf_file = request.FILES.get("pdf_file")
        ocr_option = str(request.data.get("ocr", "false")).lower() == "true"
        image_description_option = str(request.data.get("image_description", "false")).lower() == "true"
        remaining_images = int(request.data.get("rImages", 25))
        language = request.data.get("language", "English")
        batch_size = 5  

        if not pdf_file:
            return Response({"error": "PDF file is required."}, status=status.HTTP_400_BAD_REQUEST)

        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as temp_pdf:
            temp_pdf.write(pdf_file.read())
            temp_pdf_path = temp_pdf.name

        try:
            text_content = ""
            image_description_count = 0

            # Create a thread-safe counter for remaining images
            remaining_images_counter = ThreadSafeCounter(remaining_images)

            with open_pdf(temp_pdf_path) as pdf_document:
                num_pages = len(pdf_document)
                with ThreadPoolExecutor() as executor:
                    futures = []
                    for i in range(0, num_pages, batch_size):
                        batch_pages = [pdf_document[j] for j in range(i, min(i + batch_size, num_pages))]
                        for page_number, page in enumerate(batch_pages, start=i + 1):
                            futures.append(executor.submit(
                                process_page, page, page_number, ocr_option, image_description_option, language, remaining_images_counter
                            ))

                    for future in futures:
                        page_text, img_count = future.result()
                        text_content += page_text
                        image_description_count += img_count

            os.remove(temp_pdf_path)
            return Response({"text_content": text_content, "count": image_description_count}, status=status.HTTP_200_OK)

        except Exception as e:
            if os.path.exists(temp_pdf_path):
                os.remove(temp_pdf_path)
            return Response({"error": str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

def save_temporary_ppt(uploaded_file):
    with tempfile.NamedTemporaryFile(delete=False, suffix=".ppt") as temp_file:
        temp_file.write(uploaded_file.read())
        return temp_file.name


def convert_ppt_to_pptx(ppt_path):
    try:
        if not os.path.exists(ppt_path):
            raise RuntimeError(f"Input file not found: {ppt_path}")

        output_dir = os.path.dirname(ppt_path)
        libreoffice_path = "C:\\Program Files\\LibreOffice\\program\\soffice.exe" if os.name == 'nt' else "libreoffice"

        command = [
            libreoffice_path,
            "--headless",
            "--convert-to", "pptx",
            "--outdir", output_dir,
            ppt_path
        ]
        subprocess.run(command, check=True)

        pptx_path = ppt_path.replace(".ppt", ".pptx")
        if not os.path.exists(pptx_path):
            raise RuntimeError(f"Conversion failed: .pptx file not found at {pptx_path}")

        return pptx_path
    except subprocess.CalledProcessError as e:
        raise RuntimeError(f"LibreOffice conversion failed: {e}")
    except Exception as e:
        raise RuntimeError(f"Error converting PPT to PPTX: {e}")


def extract_content_from_pptx(presentation):
    slides_content = []
    try:
        for slide_index, slide in enumerate(presentation.slides, start=1):
            slide_data = {"slide_number": slide_index, "texts": [], "images": []}

            for shape in slide.shapes:
                if shape.has_text_frame:
                    slide_data["texts"].extend(paragraph.text.strip() for paragraph in shape.text_frame.paragraphs)

                if hasattr(shape, "image"):
                    image_stream = shape.image.blob
                    slide_data["images"].append(base64.b64encode(image_stream).decode("utf-8"))

            slides_content.append(slide_data)

    except Exception as e:
        raise Exception(f"Error extracting content from slides: {e}")

    return slides_content


# class PptxProcessorAPIView(APIView):
#     parser_classes = [MultiPartParser]

#     def post(self, request, *args, **kwargs):
#         pptx_file = request.FILES.get("file")
#         language = request.data.get("language", "English")
#         image_description = request.data.get("image_description", 'true')
#         remaining_images = int(request.data.get("rImages")) if request.data.get("rImages") else 25
#         print(remaining_images)
#         if not pptx_file:
#             return Response({"error": "No file uploaded."}, status=status.HTTP_400_BAD_REQUEST)

#         try:
#             image_description_count = 0
#             temp_file_path = save_temporary_ppt(pptx_file)

#             pptx_file_path = convert_ppt_to_pptx(temp_file_path) if pptx_file.name.lower().endswith(".ppt") else temp_file_path

#             presentation = Presentation(pptx_file_path)
#             slides_content = extract_content_from_pptx(presentation)

#             if str(image_description) == 'false':
#                 for slide in slides_content:
#                     del slide["images"]

#             if str(image_description) == 'true' and remaining_images > 0:
#                 for slide in slides_content:
#                     described_images = []
#                     for image_base64 in slide["images"]:
#                         if remaining_images > 0:
#                             prompt_texts = {
#                                 "English": "Describe this image in detail.",
#                                 "Arabic": "صف هذه الصورة بالتفصيل.",
#                                 "Spanish": "Describe esta imagen en detalle."
#                             }
#                             prompt_text = prompt_texts.get(language, "Describe this image in detail.")
#                             described_images.append(describe_image_with_gpt(image_base64, prompt_text))
#                             image_description_count += 1
#                             remaining_images -= 1
#                         else:
#                             break  # Stop processing if remaining_images reaches 0
#                     slide["images"] = described_images  # Only include processed images
#             return Response({"slides": slides_content, "count": image_description_count}, status=status.HTTP_200_OK)

#         except Exception as e:
#             return Response({"error": str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

class PptxProcessorAPIView(APIView):
    parser_classes = [MultiPartParser]

    def post(self, request, *args, **kwargs):
        pptx_file = request.FILES.get("file")
        language = request.data.get("language", "English")
        image_description = str(request.data.get("image_description", "true")).lower() == "true"
        remaining_images = int(request.data.get("rImages", 25))

        if not pptx_file:
            return Response({"error": "No file uploaded."}, status=status.HTTP_400_BAD_REQUEST)

        try:
            temp_file_path = save_temporary_ppt(pptx_file)
            pptx_file_path = convert_ppt_to_pptx(temp_file_path) if pptx_file.name.lower().endswith(".ppt") else temp_file_path

            presentation = Presentation(pptx_file_path)
            slides = list(presentation.slides)

            slides_content = []
            image_description_count = 0
            remaining_images_counter = ThreadSafeCounter(remaining_images)

            def extract_content_from_slide(slide):
                slide_content = {"text": "", "images": []}

                # Extract text from slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content["text"] += shape.text.strip() + "\n"

                # Extract images from slide
                for shape in slide.shapes:
                    if hasattr(shape, "image"):
                        image_stream = shape.image.blob
                        image_base64 = base64.b64encode(image_stream).decode("utf-8")
                        slide_content["images"].append(image_base64)

                return slide_content

            def process_slide(slide, slide_number):
                extracted_content = extract_content_from_slide(slide)

                if image_description and remaining_images_counter.get_value() > 0:
                    described_images = []
                    for image_base64 in extracted_content["images"]:
                        if remaining_images_counter.decrement():
                            prompt_texts = {
                                "English": "Describe this image in detail.",
                                "Arabic": "صف هذه الصورة بالتفصيل.",
                                "Spanish": "Describe esta imagen en detalle."
                            }
                            prompt_text = prompt_texts.get(language, "Describe this image in detail.")
                            described_images.append(describe_image_with_gpt(image_base64, prompt_text))
                        else:
                            break
                    extracted_content["images"] = described_images
                else:
                    extracted_content.pop("images", None)

                return extracted_content

            with ThreadPoolExecutor() as executor:
                futures = {executor.submit(process_slide, slide, i): i for i, slide in enumerate(slides)}

                for future in futures:
                    slide_content = future.result()
                    slides_content.append(slide_content)
                    image_description_count += len(slide_content.get("images", []))

            return Response({"slides": slides_content, "count": image_description_count}, status=status.HTTP_200_OK)

        except Exception as e:
            return Response({"error": str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)


@api_view(["POST"])
def create_history(request):
    serializer = serializers.HistorySerializer(data=request.data)
    if serializer.is_valid():
        serializer.save()
        return Response(serializer.data, status=status.HTTP_201_CREATED)
    return Response(serializer.errors, status=status.HTTP_400_BAD_REQUEST)


@api_view(["GET"])
def get_history(request, user_id):
    history = models.History.objects.filter(user=user_id)
    serializer = serializers.HistorySerializer(history, many=True)
    return Response(serializer.data)


@api_view(["GET"])
def get_history_by_id(request, pk):
    history = models.History.objects.get(pk=pk)
    serializer = serializers.HistorySerializer(history)
    return Response(serializer.data)

