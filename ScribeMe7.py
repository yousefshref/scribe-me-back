from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from docx import Document
import win32com.client
import pytesseract
from PIL import Image
import io
import sys
import os
import subprocess
import tempfile
import wx
import threading
import pyttsx3


def get_tesseract_path():
    try:
        base_path = sys._MEIPASS
        ocr_path = os.path.join(base_path, "Tesseract-OCR", "tesseract.exe")
        return ocr_path
    except Exception:
        return r'C:\Program Files\Tesseract-OCR\tesseract.exe'  
        
pytesseract.pytesseract.tesseract_cmd = get_tesseract_path()
        

def extract_text_from_image_using_image_path(image_path):
    try:
        image = Image.open(image_path)
        text = pytesseract.image_to_string(image)
        return text
    
    except Exception as e:
        print(f"error extracting text from image: {e}")
        return None

def extract_text_from_image(image_data):
    try:
        image = Image.open(io.BytesIO(image_data))
        text = pytesseract.image_to_string(image)
        return text
    except Exception as e:
        print(f"error extracting text from image: {e}")
        return None

def convert_from_ppt_to_pptx(file_path_ppt, file_path_pptx):
        powerpoint = win32com.client.Dispatch("Powerpoint.Application")
        presentation = powerpoint.Presentations.open(file_path_ppt)
        presentation.saveas(file_path_pptx, 24)
        presentation.Close()
        powerpoint.Quit()

def combine_text_with_different_format(paragraph):
    combine_text = ""
    current_line = None
    for run in paragraph.runs:
            if run._r.xml == current_run:
                combined_text += run.text
            else:
                combined_text += " " + run.text
                current_run = run._r.xml
                paragraph.text = combined_text
            return combined_text
    
    
def extract_text_and_images_from_pptx(file_path):
        
        text_content =[]
        slide_number = 0
        media_folder = os.path.join(os.path.dirname(file_path), 'ppt', 'media')
    
        prs = Presentation(file_path)

        for i, slide in enumerate(prs.slides):        
            slide_number += 1
            text_content.append("")
            text_content.append(f"slide {slide_number}")
            print(f"Slide {slide_number}:")

            for shapes in slide.shapes:                
                if hasattr(shapes, "text_frame"):
                    for paragraph in shapes.text_frame.paragraphs:                                
                            text = paragraph.text
                            
                            text_content.append(text)                                                                 
                            
                            

            if shapes.shape_type == MSO_SHAPE_TYPE.PICTURE:  
                                if hasattr(shapes, "image"):
                                        image_data = shapes.image.blob
                                        ocr_text = extract_text_from_image(image_data)
                                        if ocr_text:
                                                text_content.append(f"Slide {slide_number}: OCR Text from Image: {ocr_text}")

            for shape in slide.shapes: 
                if hasattr(shape, 'image'):
                    image = shape.image
                    image_filename = f"image{i}_{shape.name}."
                    image_filename += image.ext if hasattr(image, 'ext') else 'jpg'

                    image_path = os.path.join(media_folder, image_filename)
                    print(f"  Image found in slide {i + 1}")
                    print(f"    Image file: {image_path}")

                    with tempfile.NamedTemporaryFile(delete= False, suffix='.' + image.ext) as temp_file:
                        temp_filename = temp_file.name
                        temp_file.write(image.blob)

                    print(f"    Image saved to temporary file: {temp_filename}")

                    ocr_text = extract_text_from_image_using_image_path(temp_filename)

                    if ocr_text:
                        text_content.append(f"OCR text from image: {ocr_text}")

        
        return text_content 

def save_to_word_document(text_content, file_path):
    doc = Document()
    for text in text_content:
        try:
            doc.add_paragraph(text)
        except ValueError as e:
            print(f"Error adding text to the document: {e}")
            print(f"Problematic text: {text}")
    doc.save(file_path)

class TitleBar(wx.Panel):
    def __init__(self, parent):
        super().__init__(parent)
        self.SetBackgroundColour(wx.Colour(173, 216, 230))
        title_text = wx.StaticText(self, label="ScribeMe")
        title_text.SetForegroundColour(wx.WHITE)
        title_text.SetFont(wx.Font(wx.FontInfo(14).Bold()))
        
        self.SetMinSize((-1, 50))
        
        sizer = wx.BoxSizer(wx.HORIZONTAL)        
        sizer.Add(title_text, 0, wx.ALIGN_LEFT | wx.ALIGN_CENTER_VERTICAL | wx.ALL, 10)  
        self.SetSizer(sizer)
        self.Disable()
        
class MyFrame(wx.Frame):
    def __init__(self):
        super().__init__(parent=None, title='ScribeMe', size=(800, 600))
        
        
        panel = wx.Panel(self)

        title_bar = TitleBar(panel)
        
        extracted_text_label = wx.StaticText(panel, label="Extracted Text")
        extracted_text_label.SetFont(wx.Font(12, wx.FONTFAMILY_DEFAULT, wx.FONTSTYLE_NORMAL, wx.FONTWEIGHT_NORMAL))        
        
        extracted_text_label.Wrap(200)  
        extracted_text_label.SetMinSize((200, -1))  
        
        self.text_ctrl = wx.TextCtrl(panel, style=wx.TE_MULTILINE|wx.HSCROLL|wx.TE_READONLY)
        self.text_ctrl.SetValue("Please choose a powerpoint file to start text and image OCR extraction.")
        self.text_ctrl.SetFont(wx.Font(wx.FontInfo(14)))
        self.text_ctrl.SetMinSize((400, 300))
            
        
        
        
        browse_button = wx.Button(panel, label='Browse PowerPoint Files')
        browse_button.SetFont(wx.Font(12, wx.FONTFAMILY_DEFAULT, wx.FONTSTYLE_NORMAL, wx.FONTWEIGHT_BOLD))        
        browse_button.SetBackgroundColour(wx.Colour(0, 128, 0))
        browse_button.SetForegroundColour(wx.WHITE)
        browse_button.SetMinSize((200, -1))
        browse_button.Bind(wx.EVT_BUTTON, self.on_browse)
        
        accel_tbl = wx.AcceleratorTable([(wx.ACCEL_ALT, ord('B'), browse_button.GetId())])
        self.SetAcceleratorTable(accel_tbl)

        sizer = wx.BoxSizer(wx.VERTICAL)
        sizer.Add(title_bar, 0, wx.EXPAND | wx.ALL, border=5)  
        sizer.Add(extracted_text_label, flag=wx.ALIGN_LEFT | wx.ALL, border=5)  
        sizer.Add(self.text_ctrl, proportion=1, flag=wx.EXPAND | wx.ALL, border=5)
        sizer.Add(browse_button, flag=wx.ALIGN_CENTER | wx.ALL, border=5)

        self.status_bar = self.CreateStatusBar()
        
        
        panel.SetSizer(sizer)
        
        
    def on_browse(self, event):
        wildcard = "All files (*.*)|*.*"  
        dialog = wx.FileDialog(self, "Choose a file", wildcard=wildcard, style=wx.FD_OPEN | wx.FD_FILE_MUST_EXIST)
        if dialog.ShowModal() == wx.ID_OK:
            file_path = dialog.GetPath()
            if file_path.lower().endswith(('.ppt', '.pptx')):
                if file_path.lower().endswith('.ppt'):
                    temp_pptx_file_path = tempfile.mktemp(suffix='.pptx')
                    threading.Thread(target=self.extract_and_display, args=(file_path, temp_pptx_file_path)).start()
                    
                else:
                    threading.Thread(target=self.extract_and_display, args=(file_path, )).start()
                    
            else:
                wx.MessageBox("Unsupported file type. Please choose a PPT or PPTX file.", "Error", wx.OK | wx.ICON_ERROR)
                
        dialog.Destroy()
        
        self.browse_button.Enable()

    def extract_and_display(self, file_path_ppt, temp_pptx_file_path=None):
        try:
            self.update_status("Extraction started...")
            if temp_pptx_file_path:
                convert_from_ppt_to_pptx(file_path_ppt, temp_pptx_file_path)
                pptx_file_path = file_path_ppt + 'x'
                text_content = extract_text_and_images_from_pptx((temp_pptx_file_path))
                os.remove(temp_pptx_file_path)
            else:
                text_content = extract_text_and_images_from_pptx(file_path_ppt)
            
            wx.CallAfter(self.update_text_ctrl, '\n'.join(text_content))
            wx.CallAfter(self.update_status, "Extraction complete!")
            
        except Exception as e:
            wx.CallAfter(self.update_status, f"Error: {str(e)}")
        
        self.SetStatusBar(self.status_bar)
        self.update_status("Ready")
        
        
            
    def update_text_ctrl(self, text):
        self.text_ctrl.SetValue(text)
        self.status_bar.GetAccessible().SetName(text)
        self.text_ctrl.SetFocus()

    
    def update_status(self, message):
        self.status_bar.SetStatusText(message)
                
        

        
app = wx.App()
frame = MyFrame()
frame.Maximize()  
frame.Show()
app.MainLoop()

